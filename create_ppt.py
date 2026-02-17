from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def rgb_color(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = rgb_color(26, 54, 93)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = rgb_color(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = rgb_color(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_section_slide(title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = rgb_color(44, 82, 130)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = rgb_color(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb_color(26, 54, 93)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = rgb_color(255, 255, 255)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = rgb_color(50, 50, 50)
        p.space_after = Pt(12)

def add_table_slide(title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb_color(26, 54, 93)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = rgb_color(255, 255, 255)

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * table_rows)).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb_color(226, 232, 240)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

# ========== CREATE SLIDES ==========

# Slide 1: Title
add_title_slide(
    "Social Media Automation Platform",
    "Project Handover Documentation | Version 1.0.0 | February 2026"
)

# Slide 2: Agenda
add_content_slide("Agenda", [
    "Executive Summary",
    "Technology Stack",
    "System Architecture",
    "Key Features",
    "Database Schema",
    "API Endpoints",
    "Setup & Installation",
    "Deployment Guide",
    "Security & Testing"
])

# Slide 3: Executive Summary
add_section_slide("Executive Summary")

# Slide 4: Project Overview
add_content_slide("Project Overview", [
    "Comprehensive social media automation for Stage OTT content creators",
    "Multi-platform publishing: Instagram, YouTube, Facebook",
    "AI-powered captions in 4 languages (English, Hinglish, Haryanvi, Hindi)",
    "Automatic video transcoding to 6 platform-specific formats",
    "Scheduled posting with automated publishing",
    "Analytics dashboard for performance tracking"
])

# Slide 5: Key Metrics
add_table_slide("Project Metrics",
    ["Metric", "Value"],
    [
        ["Total Files", "200+"],
        ["Lines of Code", "~15,000"],
        ["Backend Routes", "20+"],
        ["Frontend Pages", "12"],
        ["React Components", "50+"],
        ["Test Coverage (Backend)", "70%"],
        ["Test Coverage (Frontend)", "60%"]
    ]
)

# Slide 6: Technology Stack
add_section_slide("Technology Stack")

# Slide 7: Backend Tech
add_table_slide("Backend Technologies",
    ["Component", "Technology", "Version"],
    [
        ["Framework", "Express.js", "4.18"],
        ["Language", "TypeScript", "5.3"],
        ["Database", "PostgreSQL", "14+"],
        ["ORM", "Prisma", "5.8"],
        ["Queue", "Redis + Bull", "7+ / 4.12"],
        ["Storage", "AWS S3", "-"],
        ["Video Processing", "FFmpeg", "-"],
        ["Authentication", "JWT", "9.0"]
    ]
)

# Slide 8: Frontend Tech
add_table_slide("Frontend Technologies",
    ["Component", "Technology", "Version"],
    [
        ["Framework", "Next.js (App Router)", "14.1"],
        ["Language", "TypeScript", "5.3"],
        ["Styling", "TailwindCSS", "3.4"],
        ["UI Components", "shadcn/ui", "-"],
        ["State Management", "Zustand", "4.4"],
        ["Server State", "React Query", "5.17"],
        ["Forms", "React Hook Form + Zod", "7.49"]
    ]
)

# Slide 9: Architecture
add_section_slide("System Architecture")

# Slide 10: Architecture Overview
add_content_slide("Architecture Overview", [
    "Frontend: Next.js 14 with App Router",
    "Reverse Proxy: Nginx with SSL/TLS & Rate Limiting",
    "Backend API: Express.js with TypeScript",
    "Database: PostgreSQL with Prisma ORM",
    "Queue System: Redis + Bull for background jobs",
    "Storage: AWS S3 for video files",
    "Workers: Video processing, Publishing, Analytics sync"
])

# Slide 11: Video Processing Pipeline
add_content_slide("Video Processing Pipeline", [
    "1. User uploads video via frontend (max 500MB)",
    "2. Original video saved to S3 (videos/raw/)",
    "3. Processing job added to Bull queue",
    "4. FFmpeg generates 6 platform-specific formats",
    "5. Thumbnails generated at 1s, 3s, 5s",
    "6. Video status updated to READY"
])

# Slide 12: Video Formats
add_table_slide("Video Output Formats",
    ["Platform", "Format", "Resolution", "Aspect Ratio"],
    [
        ["Instagram", "Reel", "1080x1920", "9:16"],
        ["Instagram", "Feed", "1080x1080", "1:1"],
        ["YouTube", "Short", "1080x1920", "9:16"],
        ["YouTube", "Video", "1920x1080", "16:9"],
        ["Facebook", "Square", "1080x1080", "1:1"],
        ["Facebook", "Landscape", "1920x1080", "16:9"]
    ]
)

# Slide 13: Features
add_section_slide("Key Features")

# Slide 14: Core Features
add_content_slide("Core Features", [
    "Video Upload & Processing - Up to 500MB with auto-transcoding",
    "Multi-Platform Publishing - Instagram, YouTube, Facebook",
    "AI Caption Generation - 4 languages, 3 variations each",
    "Scheduling System - Future posts with auto-publishing",
    "Analytics Dashboard - Views, likes, comments, engagement",
    "OAuth Integration - Secure platform authentication",
    "Token Auto-Refresh - Seamless token management"
])

# Slide 15: Frontend Pages
add_content_slide("Frontend Pages", [
    "Login/Register - User authentication",
    "Dashboard - Overview and quick actions",
    "Videos - Video library with upload",
    "Post Creator - 5-step wizard for creating posts",
    "Post Queue - Manage scheduled/published posts",
    "Calendar - Calendar view of scheduled posts",
    "Analytics - Performance metrics dashboard",
    "Accounts - Manage connected social accounts"
])

# Slide 16: Database
add_section_slide("Database Schema")

# Slide 17: Database Tables
add_content_slide("Database Tables", [
    "User - User accounts and authentication",
    "SocialAccount - Connected social media accounts",
    "Video - Uploaded videos and processed formats",
    "Post - Social media posts and scheduling",
    "Analytics - Performance metrics per post",
    "Job - Background job tracking"
])

# Slide 18: API
add_section_slide("API Documentation")

# Slide 19: API Endpoints
add_table_slide("Main API Endpoints",
    ["Method", "Endpoint", "Description"],
    [
        ["POST", "/api/auth/register", "Register new user"],
        ["POST", "/api/auth/login", "Login user"],
        ["POST", "/api/videos/upload", "Upload video (max 500MB)"],
        ["GET", "/api/videos", "List videos"],
        ["POST", "/api/posts", "Create post"],
        ["POST", "/api/posts/:id/publish", "Publish post"],
        ["POST", "/api/posts/generate-caption", "Generate AI captions"],
        ["GET", "/api/analytics", "Get analytics"]
    ]
)

# Slide 20: Rate Limits
add_table_slide("API Rate Limits",
    ["Endpoint", "Limit"],
    [
        ["General API", "100 requests / 15 minutes"],
        ["Authentication", "5 requests / 15 minutes"],
        ["Upload", "10 requests / hour"],
        ["Caption Generation", "20 requests / hour"],
        ["Analytics", "30 requests / 15 minutes"]
    ]
)

# Slide 21: Setup
add_section_slide("Setup & Installation")

# Slide 22: Prerequisites
add_content_slide("Prerequisites", [
    "Node.js 20+",
    "PostgreSQL 14+",
    "Redis 7+",
    "FFmpeg",
    "Docker & Docker Compose (recommended)",
    "AWS S3 bucket",
    "Social media API credentials (Instagram, YouTube, Facebook)"
])

# Slide 23: Quick Start
add_content_slide("Quick Start with Docker", [
    "1. Clone repository: git clone <repository-url>",
    "2. Start services: docker-compose up -d",
    "3. Install dependencies: cd backend && npm install",
    "4. Setup database: npx prisma migrate dev",
    "5. Start servers: npm run dev",
    "6. Access: Frontend (localhost:3001), Backend (localhost:3000)"
])

# Slide 24: Deployment
add_section_slide("Deployment")

# Slide 25: Deployment Steps
add_content_slide("Production Deployment", [
    "1. Configure environment: cp .env.production.example .env",
    "2. Build and deploy: docker-compose -f docker-compose.prod.yml up -d",
    "3. Setup SSL: ./deployment/scripts/setup-ssl.sh your-domain.com",
    "4. Verify: ./deployment/scripts/health-check.sh",
    "Maintenance scripts for backup, restore, logs, health checks"
])

# Slide 26: Server Requirements
add_table_slide("Server Requirements",
    ["Setup", "RAM", "CPU", "Storage", "Use Case"],
    [
        ["Minimum", "4GB", "2 cores", "50GB", "Small traffic"],
        ["Recommended", "8GB+", "4 cores", "100GB+", "High traffic"]
    ]
)

# Slide 27: Security
add_section_slide("Security & Testing")

# Slide 28: Security Features
add_content_slide("Security Measures", [
    "JWT authentication with httpOnly cookies",
    "Bcrypt password hashing (10 rounds)",
    "Helmet.js security headers (CSP, HSTS, XSS)",
    "Rate limiting per endpoint",
    "Input sanitization (XSS prevention)",
    "SQL injection prevention (Prisma ORM)",
    "HTTPS/TLS encryption",
    "Non-root Docker containers"
])

# Slide 29: Testing
add_content_slide("Testing Coverage", [
    "Backend: 70% coverage (24 unit tests, 12 integration tests)",
    "Frontend: 60% coverage (8 component tests, 4 store tests)",
    "CI/CD Pipeline: Lint > Format > Type Check > Test > Build",
    "Automated on every push/PR via GitHub Actions",
    "Security audit included in pipeline"
])

# Slide 30: Thank You
add_title_slide(
    "Thank You!",
    "Project Status: Complete & Production Ready"
)

# Save presentation
prs.save('/Users/suhani/social-media-automation/PROJECT_HANDOVER_PRESENTATION.pptx')
print("PPT created successfully!")
